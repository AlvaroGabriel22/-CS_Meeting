"""Export validation: the file exists, opens, and says what the page said."""

from __future__ import annotations

import base64
import re
from pathlib import Path

import pytest
from pypdf import PdfReader
from pptx import Presentation

PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _upload(client, path: Path, department: str = "IQC"):
    return client.post(
        "/api/uploads",
        data={"department": department, "createVersion": "true"},
        files={
            "file": (
                path.name,
                path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    ).json()


def _pdf_text(content: bytes, tmp_path: Path) -> tuple[int, str]:
    tmp_path.mkdir(parents=True, exist_ok=True)
    path = tmp_path / "export.pdf"
    path.write_bytes(content)
    reader = PdfReader(str(path))
    return len(reader.pages), "\n".join(page.extract_text() or "" for page in reader.pages)


def _deck(content: bytes, tmp_path: Path) -> Presentation:
    path = tmp_path / "export.pptx"
    path.write_bytes(content)
    return Presentation(str(path))


def _deck_text(deck: Presentation) -> str:
    return "\n".join(
        shape.text_frame.text
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )


# --------------------------------------------------------------------------- #
# 20, 22-24, 27. PDF


TITLE = "Relatorio da semana"
LINES = [
    "Rejection rate rose this month, driven by one supplier.",
    "Containment in place; 3.000 units re-inspected.",
    "Audit scheduled for next week.",
]


def _content(image_asset: int | None = None) -> dict:
    blocks = [
        {"id": f"b{index}", "type": "text", "text": line, "align": "left"}
        for index, line in enumerate(LINES)
    ]
    if image_asset is not None:
        blocks.insert(1, {"id": "img", "type": "image", "assetId": image_asset,
                          "caption": "evidencia", "align": "center", "width": 80})
    return {
        "title": TITLE,
        "columns": [{"id": "c1", "name": "Observacao"}, {"id": "c2", "name": "Acao"}],
        "rows": [
            {"id": "r1", "cells": {"c1": blocks, "c2": [
                {"id": "s1", "type": "shape", "shape": "rectangle", "color": "#1E3A5F"},
            ]}}
        ],
    }


def _with_report(client, path: Path, content: dict | None = None) -> int:
    version_id = _upload(client, path)["versionId"]
    client.put(
        f"/api/versions/{version_id}/report", json={"content": content or _content()}
    )
    return version_id


# --------------------------------------------------------------------------- #
# PDF — charts, tables, report
# --------------------------------------------------------------------------- #
def test_the_pdf_is_generated_and_opens(client, iqc_real: Path, tmp_path: Path) -> None:
    version = _upload(client, iqc_real)["versionId"]
    response = client.post(f"/api/versions/{version}/export/pdf", json={})

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/pdf"
    pages, text = _pdf_text(response.content, tmp_path)
    assert pages >= 2
    assert "IQC" in text


def test_the_pdf_states_the_version_it_came_from(client, iqc_real: Path, tmp_path: Path) -> None:
    created = _upload(client, iqc_real)
    _pages, text = _pdf_text(
        client.post(f"/api/versions/{created['versionId']}/export/pdf", json={}).content,
        tmp_path,
    )
    assert f"Version {created['versionNumber']}" in text
    assert "RawdataIQC.xlsx" in text or "Metric PPM" in text


def test_the_pdf_carries_the_three_tables_in_the_file_order(
    client, iqc_real: Path, tmp_path: Path
) -> None:
    _pages, text = _pdf_text(
        client.post(
            f"/api/versions/{_upload(client, iqc_real)['versionId']}/export/pdf", json={}
        ).content,
        tmp_path,
    )
    assert text.index("TTL") < text.index("SEC") < text.index("TNP")
    assert "Rej. Lot" in text and "Insp. Lot" in text
    # the headline row has no metric label in the file, and none is invented
    assert "\nPPM\n" not in text


def test_the_pdf_carries_the_report_the_author_built(
    client, iqc_real: Path, tmp_path: Path
) -> None:
    version = _with_report(client, iqc_real)
    _pages, text = _pdf_text(
        client.post(f"/api/versions/{version}/export/pdf", json={}).content, tmp_path
    )
    assert TITLE in text
    assert "Observacao" in text and "Acao" in text  # the columns the author named
    for line in LINES:
        assert line in text  # verbatim, block for block


def test_a_version_without_a_report_exports_without_one(
    client, iqc_real: Path, tmp_path: Path
) -> None:
    version = _upload(client, iqc_real)["versionId"]
    _pages, text = _pdf_text(
        client.post(f"/api/versions/{version}/export/pdf", json={}).content, tmp_path
    )
    assert TITLE not in text  # nothing is written on the author's behalf


def test_each_part_can_be_downloaded_on_its_own(client, iqc_real: Path, tmp_path: Path) -> None:
    """The library offers report, charts and tables separately (ADR-0038)."""
    version = _with_report(client, iqc_real)

    _pages, only_report = _pdf_text(
        client.post(
            f"/api/versions/{version}/export/pdf",
            json={"includeCharts": False, "includeTables": False},
        ).content,
        tmp_path / "report",
    )
    assert TITLE in only_report
    assert "Rej. Lot" not in only_report

    _pages, only_tables = _pdf_text(
        client.post(
            f"/api/versions/{version}/export/pdf",
            json={"includeCharts": False, "includeReport": False},
        ).content,
        tmp_path / "tables",
    )
    assert "Rej. Lot" in only_tables
    assert TITLE not in only_tables


def test_an_image_placed_in_a_cell_reaches_the_pdf(
    client, iqc_real: Path, tmp_path: Path
) -> None:
    version = _upload(client, iqc_real)["versionId"]
    uploaded = client.post(
        f"/api/versions/{version}/report/media",
        files={"file": ("evidence.png", PNG, "image/png")},
    ).json()
    client.put(
        f"/api/versions/{version}/report", json={"content": _content(uploaded["assetId"])}
    )
    response = client.post(f"/api/versions/{version}/export/pdf", json={})
    path = tmp_path / "with-image.pdf"
    path.write_bytes(response.content)
    reader = PdfReader(str(path))
    images = [image for page in reader.pages for image in page.images]
    assert images, "the evidence image is embedded"


# --------------------------------------------------------------------------- #
# PowerPoint — charts, tables, report
# --------------------------------------------------------------------------- #
def test_the_deck_is_generated_and_opens(client, iqc_real: Path, tmp_path: Path) -> None:
    version = _upload(client, iqc_real)["versionId"]
    response = client.post(f"/api/versions/{version}/export/ppt", json={})

    assert response.status_code == 200
    assert "presentationml" in response.headers["content-type"]
    deck = _deck(response.content, tmp_path)
    assert len(deck.slides) >= 4  # one chart slide + three table slides


def test_the_deck_charts_are_native_charts(client, iqc_real: Path, tmp_path: Path) -> None:
    version = _upload(client, iqc_real)["versionId"]
    deck = _deck(client.post(f"/api/versions/{version}/export/ppt", json={}).content, tmp_path)

    charts = [shape.chart for slide in deck.slides for shape in slide.shapes if shape.has_chart]
    assert len(charts) == 3, "one chart per table, side by side on one slide"
    categories = list(charts[0].plots[0].categories)
    assert categories == ["'25", "'26", "1Q", "2Q", "3Q", "Aug"]
    names = [series.name for series in charts[0].plots[0].series]
    assert names == ["SKD", "CKD", "Local", "Total"]  # parts, then the whole


def test_the_deck_tables_keep_their_structure(client, iqc_real: Path, tmp_path: Path) -> None:
    version = _upload(client, iqc_real)["versionId"]
    deck = _deck(client.post(f"/api/versions/{version}/export/ppt", json={}).content, tmp_path)

    tables = [shape.table for slide in deck.slides for shape in slide.shapes if shape.has_table]
    assert len(tables) == 3
    cells = [cell.text for table in tables for row in table.rows for cell in row.cells]
    assert "Imported" in cells and "Rej. Lot" in cells
    assert "PPM" not in cells  # never invented into the headline row
    merged = [
        cell
        for table in tables
        for row in table.rows
        for cell in row.cells
        if cell.is_merge_origin
    ]
    assert merged, "the merges of the workbook survive"


def test_the_deck_carries_the_report_as_a_native_table(
    client, iqc_real: Path, tmp_path: Path
) -> None:
    version = _with_report(client, iqc_real)
    deck = _deck(client.post(f"/api/versions/{version}/export/ppt", json={}).content, tmp_path)
    text = _deck_text(deck)

    assert TITLE in text
    tables = [shape.table for slide in deck.slides for shape in slide.shapes if shape.has_table]
    assert len(tables) == 4  # three from the workbook, one for the report
    cells = [cell.text for row in tables[-1].rows for cell in row.cells]
    assert "Observacao" in cells and "Acao" in cells
    # the blocks of one cell stay in one cell, in the author's order
    written = "\n".join(cells)
    for line in LINES:
        assert line in written
    assert written.index(LINES[0]) < written.index(LINES[1]) < written.index(LINES[2])


# --------------------------------------------------------------------------- #
# What the exports never do
# --------------------------------------------------------------------------- #
def test_neither_export_composes_a_sentence_about_the_data(
    client, iqc_real: Path, tmp_path: Path
) -> None:
    """Charts, tables and the author's words — nothing else (ADR-0033/0036)."""
    version = _upload(client, iqc_real)["versionId"]
    _pages, pdf_text = _pdf_text(
        client.post(f"/api/versions/{version}/export/pdf", json={}).content, tmp_path
    )
    deck_text = _deck_text(
        _deck(client.post(f"/api/versions/{version}/export/ppt", json={}).content, tmp_path)
    )

    forbidden = ("insight", "because", "caused", "due to", "root cause", "rose",
                 "fell", "worsening", "improving", "largest", "trend", "severity")
    for name, text in (("pdf", pdf_text), ("pptx", deck_text)):
        lowered = text.lower()
        for word in forbidden:
            assert word not in lowered, f"{name} states {word!r}"


def test_a_translated_export_keeps_every_number(client, iqc_real: Path, tmp_path: Path) -> None:
    """The report may change language; the tables and charts may not change."""
    from app.services.translation.provider import TranslationResult, register_provider

    class Prefixing:
        name = "export-fake"

        def translate(self, request):
            return TranslationResult(
                segments=[f"KO {segment}" for segment in request.segments],
                provider=self.name,
                model="fake",
            )

    import app.services.translation.service as service_module

    provider = Prefixing()
    register_provider(provider)
    original_get = service_module.get_provider
    service_module.get_provider = lambda name=None: provider
    try:
        version = _with_report(client, iqc_real)
        plain = client.post(f"/api/versions/{version}/export/pdf", json={}).content
        translated = client.post(
            f"/api/versions/{version}/export/pdf", json={"language": "ko", "translate": True}
        ).content
    finally:
        service_module.get_provider = original_get

    _a, plain_text = _pdf_text(plain, tmp_path / "plain")
    _b, translated_text = _pdf_text(translated, tmp_path / "translated")

    assert "KO " in translated_text, "the report followed the language"
    numbers = re.compile(r"\d[\d,.]*")
    assert numbers.findall(plain_text) == numbers.findall(translated_text)


def test_exporting_twice_does_not_reuse_a_stale_file(client, iqc_evolution, tmp_path: Path) -> None:
    first = _upload(client, iqc_evolution["a"])["versionId"]
    second = _upload(client, iqc_evolution["c"])["versionId"]

    _pages_a, text_a = _pdf_text(
        client.post(f"/api/versions/{first}/export/pdf", json={}).content, tmp_path / "a"
    )
    _pages_b, text_b = _pdf_text(
        client.post(f"/api/versions/{second}/export/pdf", json={}).content, tmp_path / "b"
    )
    assert text_a != text_b
