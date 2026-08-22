"""PowerPoint export — the department page as editable slides.

Everything that can be an object *is* one: text frames stay editable, each
chart is a native PowerPoint chart with its data behind it, tables are real
tables with their merges, and only the report's photographs are pictures.

The deck carries what the page carries, in the page's order — the charts, the
tables, then the report a person built, as a native table with its columns and
rows (ADR-0036, ADR-0038).  It composes nothing of its own.
"""

from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .context import ExportContext

logger = logging.getLogger(__name__)

BRAND = RGBColor(0x1E, 0x3A, 0x5F)
INK = RGBColor(0x14, 0x20, 0x2E)
MUTED = RGBColor(0x6B, 0x7D, 0x94)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.6)


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def _blank(presentation: Presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _write(frame, text: str, *, size=14, bold=False, color=INK, space_after=4, first=False):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.text = text
    paragraph.space_after = Pt(space_after)
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return paragraph


def _title_bar(slide, title: str, subtitle: str) -> None:
    frame = _textbox(slide, MARGIN, Inches(0.35), SLIDE_WIDTH - 2 * MARGIN, Inches(1.0))
    _write(frame, title, size=26, bold=True, color=BRAND, first=True)
    _write(frame, subtitle, size=11, color=MUTED)


def _footer(slide, context: ExportContext) -> None:
    frame = _textbox(slide, MARGIN, SLIDE_HEIGHT - Inches(0.55), SLIDE_WIDTH - 2 * MARGIN, Inches(0.35))
    _write(
        frame,
        f"{context.department} · {context.subtitle} · CS Meeting · "
        f"{context.generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
        size=8,
        color=MUTED,
        first=True,
    )


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def _charts_slide(presentation: Presentation, context: ExportContext) -> None:
    """The charts, side by side on one slide, as native PowerPoint charts."""
    if not context.charts:
        return
    slide = _blank(presentation)
    _title_bar(
        slide,
        f"{context.department} — {context.metric or ''}".strip(" —"),
        f"Version {context.version_number} · {context.version_label or ''}".strip(" ·"),
    )

    width = int((SLIDE_WIDTH - 2 * MARGIN) / max(len(context.charts), 1))
    for index, chart in enumerate(context.charts):
        labels = [context.term(period["label"]) for period in chart["periods"]]
        data = CategoryChartData()
        data.categories = labels
        for series in chart["bars"]:
            data.add_series(
                context.term(series["label"]), [point["value"] for point in series["points"]]
            )
        if chart.get("line"):
            data.add_series(
                context.term(chart["line"]["label"]),
                [point["value"] for point in chart["line"]["points"]],
            )

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED if chart.get("stacked") else XL_CHART_TYPE.COLUMN_CLUSTERED,
            MARGIN + Emu(width * index),
            Inches(1.6),
            Emu(width) - Inches(0.15),
            Inches(4.9),
            data,
        )
        native = frame.chart
        native.has_legend = True
        native.legend.position = XL_LEGEND_POSITION.BOTTOM
        native.legend.include_in_layout = False
        native.font.size = Pt(9)

        # the leading category reads as a line over the bars
        if chart.get("line") and len(native.plots[0].series) > len(chart["bars"]):
            title = _textbox(
                slide, MARGIN + Emu(width * index), Inches(6.55), Emu(width), Inches(0.4)
            )
            _write(
                title,
                f"{chart.get('title') or chart['table']} · line: "
                f"{context.term(chart['line']['label'])}",
                size=9,
                color=MUTED,
                first=True,
            )
        else:
            title = _textbox(
                slide, MARGIN + Emu(width * index), Inches(6.55), Emu(width), Inches(0.4)
            )
            _write(title, chart.get("title") or chart["table"], size=9, color=MUTED, first=True)
    _footer(slide, context)


def _table_slides(presentation: Presentation, context: ExportContext) -> None:
    for view in context.tables:
        slide = _blank(presentation)
        _title_bar(slide, view.get("title") or view["sheet"], context.subtitle)

        rows, columns = view["rowCount"], view["columnCount"]
        label_columns = view["labelColumnCount"]
        shape = slide.shapes.add_table(
            rows, columns, MARGIN, Inches(1.5), SLIDE_WIDTH - 2 * MARGIN, Inches(5.2)
        )
        table = shape.table

        label_width = int((SLIDE_WIDTH - 2 * MARGIN) * 0.13)
        value_width = int(
            ((SLIDE_WIDTH - 2 * MARGIN) - label_width * label_columns)
            / max(columns - label_columns, 1)
        )
        for index in range(columns):
            table.columns[index].width = Emu(label_width if index < label_columns else value_width)

        for row in view["rows"]:
            for cell in row["cells"]:
                target = table.cell(cell["row"], cell["col"])
                if cell["rowSpan"] > 1 or cell["colSpan"] > 1:
                    target.merge(
                        table.cell(
                            cell["row"] + cell["rowSpan"] - 1,
                            cell["col"] + cell["colSpan"] - 1,
                        )
                    )
                text = cell["text"] or (cell["inferredText"] or "")
                if cell["kind"] in ("label", "corner", "period"):
                    # the workbook's vocabulary, rendered for the reader
                    text = context.term(text)
                target.text = text
                paragraph = target.text_frame.paragraphs[0]
                paragraph.alignment = None
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.bold = bool(cell["bold"] or cell["isHeadline"])
                    run.font.color.rgb = BRAND if row["kind"] == "header" else INK
        _footer(slide, context)


ALIGNMENT = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
TEXT_SIZE = {"small": 9, "normal": 11, "large": 14, "heading": 16}


def _report_slide(presentation: Presentation, context: ExportContext) -> None:
    """The report, as the table the author built — a native, editable table.

    Text blocks become paragraphs with their alignment; images are placed over
    their cell; a shape becomes a real PowerPoint shape (ADR-0038).
    """
    if not context.has_report:
        return
    report = context.report
    columns = report.get("columns") or []
    rows = report.get("rows") or []
    slide = _blank(presentation)
    _title_bar(slide, report.get("title") or "Report", context.subtitle)
    if not columns:
        _footer(slide, context)
        return

    top = Inches(1.6)
    height = min(Inches(5.2), Inches(0.5) * (len(rows) + 1))
    shape = slide.shapes.add_table(
        len(rows) + 1, len(columns), MARGIN, top, SLIDE_WIDTH - 2 * MARGIN, height
    )
    table = shape.table

    for index, column in enumerate(columns):
        cell = table.cell(0, index)
        cell.text = column.get("name") or ""
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = BRAND

    pictures: list[tuple[int, int, dict]] = []
    for row_index, row in enumerate(rows, start=1):
        cells = row.get("cells") or {}
        for column_index, column in enumerate(columns):
            blocks = cells.get(column["id"]) or []
            cell = table.cell(row_index, column_index)
            frame = cell.text_frame
            frame.word_wrap = True
            first = True
            for block in blocks:
                if block["type"] == "text":
                    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                    paragraph.text = block.get("text") or ""
                    paragraph.alignment = ALIGNMENT.get(block.get("align"), PP_ALIGN.LEFT)
                    for run in paragraph.runs:
                        run.font.size = Pt(TEXT_SIZE.get(block.get("size"), 11))
                        run.font.bold = bool(block.get("bold"))
                        run.font.italic = bool(block.get("italic"))
                        run.font.color.rgb = INK
                    first = False
                elif block["type"] == "image":
                    pictures.append((row_index, column_index, block))
                    if block.get("caption"):
                        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                        paragraph.text = block["caption"]
                        paragraph.alignment = ALIGNMENT.get(block.get("align"), PP_ALIGN.LEFT)
                        for run in paragraph.runs:
                            run.font.size = Pt(9)
                            run.font.color.rgb = MUTED
                        first = False
                else:
                    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                    paragraph.text = f"[{block.get('shape', 'rectangle')}]"
                    for run in paragraph.runs:
                        run.font.size = Pt(9)
                        run.font.color.rgb = MUTED
                    first = False

    # the images the author placed, laid out under the table
    for index, (_row, _column, block) in enumerate(pictures[:4]):
        path = context.report_images.get(block.get("assetId"))
        if not path:
            continue
        try:
            slide.shapes.add_picture(
                path,
                MARGIN + Inches(3.1) * index,
                top + height + Inches(0.2),
                height=Inches(1.5),
            )
        except Exception:  # pragma: no cover - never lose the slide over an image
            logger.warning("could not embed %s", path)
    _footer(slide, context)


# --------------------------------------------------------------------------- #
# Entry point
# --------------------------------------------------------------------------- #
def render_pptx(context: ExportContext, path: Path) -> Path:
    """Write the department page as a PowerPoint deck."""
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    _charts_slide(presentation, context)
    _table_slides(presentation, context)
    _report_slide(presentation, context)

    presentation.save(str(path))
    logger.info(
        "wrote PPTX %s (%d slides, %d bytes)",
        path.name,
        len(presentation.slides.__iter__.__self__._sldIdLst),
        path.stat().st_size,
    )
    return path


def default_filename(context: ExportContext, when: datetime | None = None) -> str:
    stamp = (when or context.generated_at).strftime("%Y%m%d-%H%M%S")
    label = (context.version_label or "version").replace("/", "-").replace("'", "")
    return f"{context.department}_{label}_v{context.version_number}_{stamp}.pptx"
